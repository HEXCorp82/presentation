from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Understanding and Addressing Mental Health Issues in Young People"
subtitle.text = "Promoting Awareness and Effective Interventions\nPresenter's Name\nDate"

# Function to add a slide with title and bullet points
def add_bullet_slide(prs, title_text, subtitle_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    body = slide.placeholders[2]
    
    title.text = title_text
    subtitle.text = subtitle_text
    for point in bullet_points:
        p = body.text_frame.add_paragraph()
        p.text = point
        p.level = 0

# Slide 2: Introduction
add_bullet_slide(prs, "Introduction", "Overview of the Importance of Mental Health in Young People", [
    "Brief overview of the importance of mental health in young people.",
    "Mental health is crucial for overall well-being.",
    "Increasing prevalence of mental health issues among youth."
])

# Slide 3: Defining Mental Health
add_bullet_slide(prs, "Defining Mental Health", "Understanding the Scope of Mental Health", [
    "Mental health includes emotional, psychological, and social well-being.",
    "Affects how individuals think, feel, and act."
])

# Slide 4: Common Mental Health Issues in Young People
add_bullet_slide(prs, "Common Mental Health Issues in Young People", "Recognizing Key Challenges", [
    "Anxiety disorders",
    "Depression",
    "Attention-Deficit/Hyperactivity Disorder (ADHD)",
    "Eating disorders",
    "Substance abuse"
])

# Slide 5: Signs and Symptoms
add_bullet_slide(prs, "Signs and Symptoms", "Identifying Indications of Mental Health Issues", [
    "Changes in mood and behavior",
    "Withdrawal from social interactions",
    "Decline in academic performance",
    "Physical symptoms (e.g., headaches, stomachaches)"
])

# Slide 6: Causes and Risk Factors
add_bullet_slide(prs, "Causes and Risk Factors", "Exploring Influential Factors", [
    "Genetic predisposition",
    "Environmental factors (e.g., family issues, bullying)",
    "Social media and technology impact",
    "Academic pressure"
])

# Slide 7: The Impact of Mental Health Issues
add_bullet_slide(prs, "The Impact of Mental Health Issues", "Understanding Consequences", [
    "Academic performance",
    "Social relationships",
    "Physical health",
    "Long-term consequences"
])

# Slide 8: Breaking the Stigma
add_bullet_slide(prs, "Breaking the Stigma", "Promoting Openness and Acceptance", [
    "Encourage open discussions",
    "Educate about mental health",
    "Promote understanding and empathy"
])

# Slide 9: Approaches to Address Mental Health Issues
add_bullet_slide(prs, "Approaches to Address Mental Health Issues", "Effective Intervention Strategies", [
    "Counseling and therapy",
    "Medication management",
    "School-based programs",
    "Family support and involvement"
])

# Slide 10: Role of Parents and Guardians
add_bullet_slide(prs, "Role of Parents and Guardians", "Supporting Mental Health at Home", [
    "Recognize signs and symptoms",
    "Provide a supportive environment",
    "Encourage professional help",
    "Open communication"
])

# Slide 11: Role of Schools and Educators
add_bullet_slide(prs, "Role of Schools and Educators", "Creating Supportive Learning Environments", [
    "Implement mental health programs",
    "Train staff to recognize and address issues",
    "Provide resources and support"
])

# Slide 12: Community and Policy Support
add_bullet_slide(prs, "Community and Policy Support", "Fostering Collective Responsibility", [
    "Community programs and resources",
    "Advocacy for mental health policies",
    "Collaboration with healthcare providers"
])

# Slide 13: Case Studies and Success Stories
add_bullet_slide(prs, "Case Studies and Success Stories", "Illustrating Effective Solutions", [
    "Case study 1: School-based mental health program",
    "Case study 2: Community support initiative"
])

# Slide 14: Resources for Further Help
add_bullet_slide(prs, "Resources for Further Help", "Providing Accessible Support Channels", [
    "Mental health hotlines",
    "Websites and online communities",
    "Local mental health clinics and professionals"
])

# Slide 15: Conclusion
add_bullet_slide(prs, "Conclusion", "Emphasizing Continuous Support", [
    "Recap importance of addressing mental health in youth",
    "Encourage proactive measures",
    "Promote ongoing education and support"
])

# Slide 16: Q&A
add_bullet_slide(prs, "Q&A", "Engaging with the Audience", [
    "Encourage audience to ask questions",
    "Provide thoughtful and informative answers"
])

# Save the presentation
prs.save("Understanding_and_Addressing_Mental_Health_Issues_in_Young_People.pptx")
